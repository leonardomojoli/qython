# qython/backend/services/academic_services/material_generation_service.py

import os
import json
import re
import uuid
from typing import Dict, Any, Tuple, List, Optional
from dotenv import load_dotenv
import logging
from pptx import Presentation
from pptx.util import Inches
from ..llm_services import client, PRIMARY_LLM_MODEL, FALLBACK_LLM_MODEL, MAX_RETRIES, INITIAL_BACKOFF, MAX_BACKOFF, _get_thinking_config_for_model
from google.genai import types
import time
from ...config import Config
from ...services import stock_image_service
# demjson3 removed for security (accepts Python expressions from untrusted LLM output)
import shutil
from fpdf import FPDF

load_dotenv()

logger = logging.getLogger("qython_logger")

# Schema de saída estruturada para questionnaire_objective.
# O Gemini usa constrained decoding p/ garantir JSON sintaticamente válido —
# previne "Expecting ',' delimiter" e outros erros de sintaxe do lite+thinking.
_QUESTIONNAIRE_OBJECTIVE_SCHEMA = types.Schema(
    type=types.Type.OBJECT,
    required=['questionario_objetivo'],
    properties={
        'questionario_objetivo': types.Schema(
            type=types.Type.ARRAY,
            items=types.Schema(
                type=types.Type.OBJECT,
                required=['pergunta', 'alternativas', 'resposta_correta', 'justificativa', 'dificuldade', 'topico'],
                properties={
                    'pergunta':          types.Schema(type=types.Type.STRING),
                    'alternativas':      types.Schema(
                        type=types.Type.ARRAY,
                        items=types.Schema(type=types.Type.STRING)
                    ),
                    'resposta_correta':  types.Schema(type=types.Type.STRING, enum=['a', 'b', 'c', 'd', 'e', 'f']),
                    'justificativa':     types.Schema(type=types.Type.STRING),
                    'dificuldade':       types.Schema(type=types.Type.STRING, enum=['facil', 'medio', 'dificil']),
                    'topico':            types.Schema(type=types.Type.STRING),
                    # rótulo do texto-base ao qual a questão se ancora ("Texto I").
                    # ⚠️ Sem declarar aqui, o constrained decoding IMPEDE o modelo de
                    # emitir o campo — foi o que aconteceu: ele escrevia "de acordo com
                    # o texto" sem conseguir declarar texto nenhum.
                    'texto_base':        types.Schema(type=types.Type.STRING),
                }
            )
        ),
        # Textos de apoio compartilhados por 2-4 questões (formato clássico de banca).
        'textos_base': types.Schema(
            type=types.Type.ARRAY,
            items=types.Schema(
                type=types.Type.OBJECT,
                required=['rotulo', 'conteudo'],
                properties={
                    'rotulo':   types.Schema(type=types.Type.STRING),
                    'conteudo': types.Schema(type=types.Type.STRING),
                    'fonte':    types.Schema(type=types.Type.STRING),
                }
            )
        ),
    }
)

def _schema_questionario_objetivo(com_texto_base: bool):
    """Schema do questionário objetivo, com ou sem `textos_base`.

    ⚠️ Texto-base ("Texto I" ancorando várias questões) é formato de PROVA DE BANCA
    (Meus Concursos). O Produtor de Materiais gera questionário de ESTUDO a partir do
    material do usuário — ali esse formato é intrusão, não recurso. Como a geração usa
    constrained decoding, tirar o campo do schema é garantia dura: o modelo não consegue
    emiti-lo nem que o prompt mandasse. Foi assim que a regra vazou de um produto para o
    outro — prompt e schema eram compartilhados."""
    props = dict(_QUESTIONNAIRE_OBJECTIVE_SCHEMA.properties)
    if not com_texto_base:
        props.pop('textos_base', None)
        item = props['questionario_objetivo'].items
        item_props = dict(item.properties)
        item_props.pop('texto_base', None)
        props['questionario_objetivo'] = types.Schema(
            type=types.Type.ARRAY,
            items=types.Schema(type=types.Type.OBJECT, required=item.required, properties=item_props),
        )
    return types.Schema(type=types.Type.OBJECT, required=['questionario_objetivo'], properties=props)


# New character-based weight limit for a slide. This is more reliable than an abstract complexity score.
# This value was determined empirically to be a good balance for readability.
MAX_CHAR_WEIGHT_PER_SLIDE = 1000

def _get_content_block_weight(block: Dict) -> int:
    """Calculates the character 'weight' of a content block."""
    block_type = block.get("type")
    weight = 0
    
    if block_type in ["text", "key_takeaway"] and block.get("points"):
        weight = sum(len(str(point)) for point in block.get("points", []))
    elif block_type == "table" and block.get("rows"):
        # Table weight is the sum of characters in all cells, with a multiplier for more columns
        num_cols = len(block.get("columns", []))
        char_count = sum(len(str(cell)) for row in block.get("rows", []) for cell in row)
        # A simple heuristic: more columns take up proportionally more vertical space.
        weight = char_count * (1 + num_cols * 0.1)
    elif block_type == "clinical_vignette":
        weight = len(block.get("scenario", "")) + len(block.get("question", "")) + len(block.get("answer", ""))
    elif block_type == "image_suggestion":
        # Assign a fixed, significant weight to images to ensure they get enough space
        weight = 400
        
    return int(weight)

def _post_process_slides_for_size(slides: List[Dict]) -> List[Dict]:
    """
    Analyzes slides and splits them based on a character weight limit to ensure
    content fits well on each slide.
    """
    processed_slides = []
    for slide in slides:
        original_title = slide.get("title", "Slide")
        content_blocks = slide.get("content", [])
        
        if not content_blocks:
            processed_slides.append(slide)
            continue

        current_slide_blocks = []
        current_weight = 0
        part_counter = 1

        for block in content_blocks:
            block_weight = _get_content_block_weight(block)

            # --- Intelligent Splitting for Text Blocks ---
            # If the block itself is a text/list and is too large, split it internally.
            if block.get("type") in ["text", "key_takeaway"] and block_weight > MAX_CHAR_WEIGHT_PER_SLIDE:
                points = block.get("points", [])
                current_points_block = {"type": block["type"], "points": []}
                current_points_weight = 0

                for point in points:
                    point_weight = len(str(point))
                    if current_weight + current_points_weight + point_weight > MAX_CHAR_WEIGHT_PER_SLIDE and current_slide_blocks:
                        # Finalize the current slide before adding the oversized point
                        if current_points_block["points"]:
                            current_slide_blocks.append(current_points_block)
                        
                        new_title = f"{original_title} (Parte {part_counter})"
                        processed_slides.append({"title": new_title, "content": current_slide_blocks})
                        part_counter += 1
                        
                        # Start new slide with the current point
                        current_slide_blocks = []
                        current_weight = 0
                        current_points_block = {"type": block["type"], "points": [point]}
                        current_points_weight = point_weight
                    else:
                        current_points_block["points"].append(point)
                        current_points_weight += point_weight
                
                # Add the remaining points block to the current slide
                if current_points_block["points"]:
                    current_slide_blocks.append(current_points_block)
                    current_weight += current_points_weight

            # --- Standard Block Splitting ---
            # If adding the next block (which is not a splittable list) exceeds the limit
            elif current_weight + block_weight > MAX_CHAR_WEIGHT_PER_SLIDE and current_slide_blocks:
                # Finalize the current slide
                new_title = f"{original_title} (Parte {part_counter})"
                processed_slides.append({"title": new_title, "content": current_slide_blocks})
                part_counter += 1
                
                # Start a new slide with the current block
                current_slide_blocks = [block]
                current_weight = block_weight
            else:
                # Add the block to the current slide
                current_slide_blocks.append(block)
                current_weight += block_weight

        # Add the last remaining slide if it has content
        if current_slide_blocks:
            is_multipart = part_counter > 1 or any(s['title'].startswith(f"{original_title} (Parte") for s in processed_slides)
            new_title = f"{original_title} (Parte {part_counter})" if is_multipart else original_title
            processed_slides.append({"title": new_title, "content": current_slide_blocks})

    return processed_slides

def create_presentation_from_json(data: dict, user_id: int) -> str:
    prs = Presentation()
    # Set slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = data.get("title", "Apresentação Gerada por Qython")
    subtitle_shape.text = f"Conteúdo sobre: {data.get('theme', 'Tópico Específico')}"
    
    sized_slides = _post_process_slides_for_size(data.get("slides", []))
    content_slide_layout = prs.slide_layouts[5] # Use a blank layout

    for slide_data in sized_slides:
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add title to the slide
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_shape.text_frame.text = slide_data.get("title", "Slide sem Título")

        # --- DYNAMIC LAYOUT LOGIC ---
        # Only consider blocks that have content (i.e., successful images)
        content_blocks = [
            block for block in slide_data.get("content", []) 
            if block.get("type") != 'image_suggestion' or 
               (block.get("type") == 'image_suggestion' and block.get("temp_image_path") and os.path.exists(block.get("temp_image_path")))
        ]
        
        text_blocks = [b for b in content_blocks if b.get("type") != 'image_suggestion']
        image_blocks = [b for b in content_blocks if b.get("type") == 'image_suggestion']

        if len(text_blocks) == 1 and len(image_blocks) == 1:
            # Side-by-side layout
            _add_content_block(slide, text_blocks[0], {'left': 0.5, 'top': 1.2, 'width': 4.5, 'height': 4.2})
            _add_content_block(slide, image_blocks[0], {'left': 5.2, 'top': 1.5, 'width': 4.3, 'height': 3.5})
        elif len(image_blocks) >= 1 and not text_blocks:
            # Image-only layout (center the first image)
            _add_content_block(slide, image_blocks[0], {'left': 2, 'top': 1.2, 'width': 6, 'height': 4})
        else:
            # Default to vertical stacking for all other cases
            current_y = 1.2
            total_height_available = 4.2 # Total inches available for content
            
            for block in content_blocks:
                # A simple improvement: give more space to blocks with more points/rows
                num_items = len(block.get("points", [])) or len(block.get("rows", [])) or 1
                proportional_height = (num_items / 10) * total_height_available # Adjust '10' as needed
                height = max(0.5, min(proportional_height, total_height_available - (current_y - 1.2))) # Ensure min/max height

                _add_content_block(slide, block, {'left': 0.5, 'top': current_y, 'width': 9, 'height': height})
                current_y += height + 0.1

    output_dir = Config.SLIDESHOW_FOLDER
    os.makedirs(output_dir, exist_ok=True)
    filename = f"slideshow_{user_id}_{uuid.uuid4()}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    logger.info(f"Rich content presentation saved at: {filepath}")
    relative_path = os.path.join(Config.STATIC_URL_PATH_PREFIX.strip('/'), 'uploads', 'slideshows', filename)
    return relative_path.replace("\\", "/")

def _add_content_block(slide, block, layout_info):
    """Adds a single content block to a slide using dynamic layout info."""
    block_type = block.get("type")
    
    # Use provided layout dimensions
    left, top, width, height = (Inches(layout_info.get(key, 1.0)) for key in ['left', 'top', 'width', 'height'])

    if block_type == "text" or (block_type == "key_takeaway" and block.get("points")):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for item in block.get("points", []):
            p = tf.add_paragraph()
            p.text = item.lstrip()
            p.level = 1 if item.strip().startswith("  ") else 0

    elif block_type == "table" and block.get('rows') and block.get('columns'):
        rows = len(block['rows']) + 1
        cols = len(block['columns'])
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        for i, col_name in enumerate(block['columns']):
            table.cell(0, i).text = col_name
            table.cell(0, i).text_frame.paragraphs[0].font.bold = True
        for r, row_data in enumerate(block['rows']):
            for c, cell_data in enumerate(row_data):
                table.cell(r + 1, c).text = str(cell_data)

    elif block_type == "image_suggestion":
        image_path = block.get("temp_image_path")
        if image_path and os.path.exists(image_path):
            # For images, we use height to maintain aspect ratio
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            logger.warning(f"Skipping image block in PPTX for failed prompt: {block.get('description')}")

    elif block_type == "clinical_vignette":
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p_title = tf.add_paragraph(); p_title.text = block.get('title', 'Caso Clínico'); p_title.font.bold = True
        
        p_scenario_header = tf.add_paragraph(); p_scenario_header.text = "Cenário:"; p_scenario_header.font.italic = True
        p_scenario = tf.add_paragraph(); p_scenario.text = block.get('scenario', ''); p_scenario.level = 1
        p_question_header = tf.add_paragraph(); p_question_header.text = "Pergunta:"; p_question_header.font.italic = True
        p_question = tf.add_paragraph(); p_question.text = block.get('question', ''); p_question.level = 1
        p_answer_header = tf.add_paragraph(); p_answer_header.text = "Resposta Esperada:"; p_answer_header.font.italic = True
        p_answer = tf.add_paragraph(); p_answer.text = block.get('answer', ''); p_answer.level = 1

# Create a custom PDF class to handle headers and footers
class PDF(FPDF):
    def header(self):
        # Find the logo path relative to the current file
        # This assumes a specific directory structure, adjust if needed
        # services/academic_services/material_generation_service.py -> backend/ -> project_root/
        try:
            current_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(os.path.dirname(os.path.dirname(current_dir)))
            logo_path = os.path.join(project_root, 'packages', 'web', 'src', 'assets', 'qython-imagotipo.png')
            
            if os.path.exists(logo_path):
                # Add logo to the top left
                self.image(logo_path, 10, 8, 25)
                self.ln(10) # Add some space below the logo
        except Exception as e:
            logger.warning(f"Could not find or place the logo in the PDF header: {e}")

    def footer(self):
        # Position at 1.5 cm from bottom
        self.set_y(-15)
        self.set_font('DejaVu', 'I', 8) # Use Italic for the footer
        
        # Footer text
        footer_text = f"Gerado por Qython, seu assistente médico inteligente | {Config.WEB_BASE_URL} | Página {self.page_no()}"
        self.cell(0, 10, footer_text, 0, 0, 'C')


def create_pdf_from_json(data: dict, user_id: int) -> str:
    """Generates a professional, continuous PDF document from slideshow JSON data."""
    
    # Use our new custom PDF class
    pdf = PDF()
    
    # Register all font styles to prevent crashes with bold/italic text
    font_dir = os.path.dirname(Config.PDF_FONT_PATH)
    pdf.add_font('DejaVu', '', Config.PDF_FONT_PATH, uni=True)
    pdf.add_font('DejaVu', 'B', os.path.join(font_dir, 'DejaVuSans-Bold.ttf'), uni=True)
    pdf.add_font('DejaVu', 'I', os.path.join(font_dir, 'DejaVuSans-Oblique.ttf'), uni=True)
    pdf.add_font('DejaVu', 'BI', os.path.join(font_dir, 'DejaVuSans-BoldOblique.ttf'), uni=True)

    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # --- Title Page ---
    pdf.set_font('DejaVu', 'B', 24) # Bold Title
    pdf.multi_cell(0, 15, data.get("title", "Apresentação"), align='C', ln=True)
    pdf.ln(10)
    pdf.set_font('DejaVu', '', 16)
    pdf.multi_cell(0, 10, data.get("theme", ""), align='C', ln=True)

    # --- Content Flow ---
    # We no longer add a new page for every slide
    for slide_data in data.get("slides", []):
        # Add a page break before each new slide for clear separation, but not for the very first one
        if pdf.page_no() > 1 or pdf.get_y() > 50: # Add page break if not at the top of a new page
             pdf.add_page()

        pdf.set_font('DejaVu', 'B', 16) # Bold slide titles
        pdf.multi_cell(0, 10, slide_data.get("title", "Slide sem Título"), ln=True, border='B')
        pdf.ln(8)

        pdf.set_font('DejaVu', '', 10)
        for block in slide_data.get("content", []):
            # Check if there's enough space for the next block, add a page if not
            # This is a simple heuristic, can be improved if needed
            if pdf.get_y() > pdf.h - 40: # 40mm from bottom
                pdf.add_page()

            block_type = block.get("type")
            if block_type == "text":
                for point in block.get("points", []):
                    pdf.multi_cell(0, 5, f"  •  {point.strip()}", ln=True)
                pdf.ln(5)

            elif block_type == "key_takeaway" and block.get("points"):
                pdf.set_fill_color(230, 247, 255)
                pdf.set_text_color(0, 102, 153)
                pdf.set_font('DejaVu', 'B', 12) # Bold
                pdf.multi_cell(0, 8, "Pontos-Chave", ln=True, fill=True, border=1)
                pdf.set_font('DejaVu', '', 10)
                pdf.set_text_color(0, 0, 0)
                for point in block.get("points", []):
                    pdf.multi_cell(0, 5, f"  •  {point.strip()}", ln=True)
                pdf.ln(5)

            elif block_type == "clinical_vignette":
                pdf.set_font('DejaVu', 'B', 11)
                pdf.multi_cell(0, 6, f"Caso Clínico: {block.get('title', '')}", ln=True)
                pdf.set_font('DejaVu', '', 10)
                pdf.multi_cell(0, 5, f"Cenário: {block.get('scenario', '')}", ln=True)
                pdf.multi_cell(0, 5, f"Pergunta: {block.get('question', '')}", ln=True)
                pdf.multi_cell(0, 5, f"Resposta: {block.get('answer', '')}", ln=True)
                pdf.ln(5)

            elif block_type == "table" and block.get("columns") and block.get("rows"):
                pdf.set_font('DejaVu', 'B', 11)
                pdf.multi_cell(0, 6, block.get('title', 'Tabela'), ln=True)
                pdf.set_font('DejaVu', '', 9)
                
                with pdf.table(col_widths=(pdf.w - 20) / len(block["columns"])) as table:
                    header = table.row()
                    for col_name in block["columns"]:
                        header.cell(col_name)
                    
                    for row_data in block["rows"]:
                        row = table.row()
                        for cell_data in row_data:
                            row.cell(str(cell_data))
                pdf.ln(5)

            elif block_type == "image_suggestion" and block.get("temp_image_path"):
                image_path = block["temp_image_path"]
                if os.path.exists(image_path):
                    pdf.image(image_path, w=pdf.w - 40, x=20, h=0)
                    pdf.ln(5)

    # --- Finalization ---
    output_dir = os.path.join(Config.PERMANENT_UPLOAD_FOLDER, 'temp_pdfs')
    os.makedirs(output_dir, exist_ok=True)
    filename = f"document_pdf_{user_id}_{uuid.uuid4()}.pdf"
    filepath = os.path.join(output_dir, filename)
    pdf.output(filepath)
    
    logger.info(f"PDF document saved at: {filepath}")
    return filepath

# --- UTILITY AND PROMPT FUNCTIONS ---

def extract_json_from_response(response: str) -> dict:
    """
    Extrai um objeto JSON da resposta do modelo, TOLERANTE a dados extras DEPOIS do
    objeto. Modelos com thinking (ex.: lite + thinking_level=HIGH) às vezes anexam
    texto ou um 2º objeto após o JSON válido — e o json.loads estrito quebrava com
    "Extra data". Usa raw_decode p/ pegar o PRIMEIRO valor JSON completo e ignorar a cauda.
    """
    logger.debug(f"Attempting to extract JSON from response (first 300 chars): {response[:300]}")

    # First, find the JSON block, whether it's in a code block or not
    match = re.search(r'```json\s*(.*?)\s*```', response, re.DOTALL)
    if match:
        candidate = match.group(1).strip()
    else:
        start_index = response.find('{')
        if start_index == -1:
            raise ValueError("No valid JSON object found in the response.")
        candidate = response[start_index:]

    decoder = json.JSONDecoder()
    try:
        # Pega o 1º objeto JSON completo e ignora qualquer coisa depois (dados extras).
        obj, _end = decoder.raw_decode(candidate)
        return obj
    except json.JSONDecodeError as e:
        # Fallback: recorte clássico { ... último } (caso haja prefixo inesperado),
        # ainda via raw_decode p/ tolerar a cauda.
        s_idx = response.find('{')
        e_idx = response.rfind('}')
        if s_idx != -1 and e_idx > s_idx:
            try:
                obj, _end = decoder.raw_decode(response[s_idx:e_idx + 1])
                return obj
            except json.JSONDecodeError:
                pass
        # Fallback 2: json-repair — corrige erros comuns de LLM: vírgulas faltando,
        # colchetes não fechados, aspas não escapadas, etc. Mais seguro que demjson3
        # (não aceita expressões Python, só repara sintaxe JSON).
        try:
            import json_repair
            repaired = json_repair.repair_json(candidate, return_objects=True)
            if isinstance(repaired, (dict, list)):
                logger.warning("[JSON] json_repair recuperou JSON malformado do LLM")
                return repaired
        except Exception:
            pass
        logger.error(f"json parse failed: {e}. Content: {candidate[:500]}...")
        raise ValueError(f"Failed to parse JSON from the model's response: {e}")

def parse_text_to_quiz_json(text_response: str) -> dict:
    """Converts a structured text string into a quiz JSON object."""
    questions = []
    for block in text_response.strip().split('---'):
        if not block.strip(): continue
        q_data = {}
        lines = block.strip().split('\n')
        for line in lines:
            if ':' in line:
                key, value = line.split(':', 1)
                key = key.strip().upper()
                value = value.strip()
                if key == 'P': q_data['pergunta'] = value
                elif key in ('A', 'B', 'C', 'D', 'E', 'F'): q_data.setdefault('alternativas', []).append(value)
                elif key == 'R': q_data['resposta_correta'] = value.lower()
                elif key == 'J': q_data['justificativa'] = value
        if q_data.get('pergunta'):
            questions.append(q_data)
    return {"questionario_objetivo": questions}

def _shuffle_quiz_positions(content: dict) -> None:
    """Embaralha a ORDEM das alternativas de cada questão objetiva (remapeando resposta_correta),
    in-place. O modelo concentra a correta em B/C/D e evita A/E (viés reportado em feedback do
    usuário); embaralhar garante posição UNIFORME sem alterar o conteúdo das alternativas."""
    import random
    if not isinstance(content, dict):
        return
    questions = content.get('questionario_objetivo')
    if not isinstance(questions, list):
        return
    for q in questions:
        if not isinstance(q, dict):
            continue
        alts = q.get('alternativas')
        corr = q.get('resposta_correta')
        if not isinstance(alts, list) or len(alts) < 2 or not isinstance(corr, str) or not corr.strip():
            continue
        idx = ord(corr.strip().lower()[0]) - ord('a')
        if not (0 <= idx < len(alts)):
            continue
        order = list(range(len(alts)))
        random.shuffle(order)
        q['alternativas'] = [alts[i] for i in order]
        q['resposta_correta'] = chr(ord('a') + order.index(idx))


def _mapa_de_cobertura(prior_questions):
    """Histograma de tópicos já cobertos, extraído da própria avoid-list.

    Pedir ao modelo que compare cada questão nova contra 100 enunciados é tarefa cara —
    ele escreve as questões e "esquece" a lista (medido: 39 pares idênticos entre
    questionários da mesma biblioteca, com a lista inteira no prompt). Um mapa de
    saturação é acionável: em vez de "não repita estas 100", ele lê "estes assuntos já
    têm 12 questões, estes têm zero" e sabe para onde ir ANTES de escrever a primeira.
    """
    from collections import Counter
    topicos = Counter()
    for linha in prior_questions:
        m = re.search(r'·\s*([a-z_]+)\]\s*$', str(linha))
        if m:
            topicos[m.group(1)] += 1
    if not topicos:
        return ""
    saturados = [f"{t} ({n})" for t, n in topicos.most_common() if n >= 3]
    poucos = [f"{t} ({n})" for t, n in topicos.most_common() if n < 3]
    partes = []
    if saturados:
        partes.append("JÁ SATURADOS (evite): " + ", ".join(saturados[:12]))
    if poucos:
        partes.append("POUCO EXPLORADOS (prefira): " + ", ".join(poucos[:12]))
    return "\n".join(partes)


def _avoid_repetition_block(prior_questions):
    """Bloco de prompt listando questões JÁ geradas (avoid-list). Vazio se não houver."""
    if not prior_questions:
        return ""
    cobertura = _mapa_de_cobertura(prior_questions)
    listed = "\n".join(f"- {q}" for q in prior_questions)
    mapa = f"\n**MAPA DE COBERTURA deste material** (use isto para ESCOLHER os assuntos antes de escrever):\n{cobertura}\n" if cobertura else ""
    return f"""
**↻ VARIE — o que JÁ FOI COBRADO nas provas anteriores (não repita):**
{mapa}
⚠️ **PRIMEIRO PASSO, antes de escrever qualquer questão:** percorra o material-fonte e
escolha os assuntos que você vai cobrir, priorizando o que está POUCO EXPLORADO acima.
Só depois escreva as questões. Escrever primeiro e conferir a lista depois é o que faz
sair questão repetida.
⛔ **PROIBIDO reaproveitar enunciado**: nenhuma questão pode ter o mesmo enunciado de uma
da lista, nem uma versão com palavras trocadas. Isso não é "revisitar um ponto de alto
rendimento" — é entregar a mesma questão duas vezes, e invalida a prova para quem estuda.
Cada linha traz o enunciado e, entre colchetes, o NÚCLEO que ele cobrou (o gabarito e o
tópico). Leia como "este ponto já está cobrado", não como "esta frase já foi usada":
reescrever a pergunta e manter o mesmo núcleo CONTINUA sendo repetição.
Regras:
1. A grande maioria das questões deve abrir pontos/subtemas que NÃO aparecem na lista.
2. No máximo uma ou duas podem revisitar um ponto de alto rendimento — e só por outro
   ÂNGULO (outro aspecto do tema, outro cenário, outra competência exigida), nunca a
   mesma pergunta com palavras, idade ou números trocados.
3. Varie também o MOLDE do comando. Se a lista já tem "No Texto I, o termo X estabelece
   uma relação de sentido de:", não entregue a mesma fôrma com outro termo — troque a
   tarefa (reescrever mantendo o sentido, identificar função sintática, justificar
   pontuação, inferir intenção do autor).
4. Se o material-fonte deste bloco for pequeno e os pontos óbvios já estiverem na lista,
   vá para o detalhe: exceções, critérios numéricos, comparações, casos-limite, aspectos
   secundários do mesmo assunto — em vez de refazer o item que já existe.
{listed}
"""


def _banca_profile_block(banca_profile):
    """Bloco de prompt com o perfil pesquisado da banca/prova (dossiê CONFIRMADO do card
    "Meus Concursos") e, quando anexadas, as PROVAS ANTERIORES. Vazio se não houver.
    Guia FORMATO/ESTILO/ênfase de temas E serve de inspiração de conteúdo (banca repete)."""
    if not banca_profile or not str(banca_profile).strip():
        return ""
    return f"""**📋 PERFIL DA BANCA/PROVA (descrição do usuário, pesquisa confirmada e/ou provas anteriores) — IMITE este estilo:**
As questões devem refletir o FORMATO, o ESTILO de cobrança e os TEMAS típicos desta prova,
descritos abaixo. Priorize os assuntos mais cobrados e o jeito de a banca formular. Se este
bloco incluir PROVAS ANTERIORES, os temas cobrados nelas também são fonte legítima de
INSPIRAÇÃO de conteúdo (bancas repetem assuntos e até questões) — reformule com novo
enunciado, ângulo e alternativas; NUNCA copie uma questão literalmente. Mantenha o
EQUILÍBRIO: a maioria das questões nasce do material-fonte, varrendo o programa amplamente;
as provas antigas inspiram uma minoria, espalhada por temas diferentes. Se o perfil divergir
do material-fonte, o perfil manda no ESTILO/ênfase; a base factual segue o material.
\"\"\"
{str(banca_profile).strip()}
\"\"\"

"""


def generate_objective_quiz_prompt(subject_context, prior_questions=None, question_count=25,
                                   num_alternatives=5, for_exam: bool = False):
    """Monta o prompt do questionário objetivo do produto certo.

    ⚠️ Os dois prompts vivem em `quiz_prompts.py` e são SEPARADOS de propósito: questionário
    de ESTUDO (Produtor de Materiais) e SIMULAÇÃO de prova de banca (Meus Concursos) são
    produtos diferentes. Enquanto compartilharam um construtor com flags, regra de um vazou
    para o outro duas vezes no mesmo dia. As regras de OFÍCIO da questão continuam
    compartilhadas lá dentro — essas precisam mesmo de fonte única."""
    from . import quiz_prompts

    avoid_block = _avoid_repetition_block(prior_questions)
    n = question_count or 25
    na = max(2, min(6, int(num_alternatives or 5)))
    _letters = ['A', 'B', 'C', 'D', 'E', 'F'][:na]
    _lowers = [c.lower() for c in _letters]
    lowers_disp = ', '.join(_lowers)                                              # "a, b, c, d, e"
    lowers_or = (', '.join(_lowers[:-1]) + ' ou ' + _lowers[-1]) if na > 1 else _lowers[0]  # "a, b, c, d ou e"

    # Exemplos JSON com EXATAMENTE `na` alternativas (não ancora o modelo em 5).
    def _build_example(correct, distractors, want_pos):
        pos = min(want_pos, na - 1)
        alts = list(distractors[:na - 1])
        alts.insert(pos, correct)
        return ',\n        '.join(f'"{a}"' for a in alts), chr(ord('a') + pos)
    ex1_alts, ex1_resp = _build_example(
        "Streptococcus pneumoniae",
        ["Mycoplasma pneumoniae", "Haemophilus influenzae", "Legionella pneumophila", "Vírus Influenza", "Klebsiella pneumoniae"],
        2)
    ex2_alts, ex2_resp = _build_example(
        "Infarto agudo do miocárdio",
        ["Embolia pulmonar", "Dissecção aórtica", "Pericardite aguda", "Refluxo gastroesofágico", "Espasmo esofágico"],
        0)
    ctx = {
        'n': n, 'na': na, 'lowers_disp': lowers_disp, 'lowers_or': lowers_or,
        'ex1_alts': ex1_alts, 'ex1_resp': ex1_resp, 'ex2_alts': ex2_alts, 'ex2_resp': ex2_resp,
        'subject_context': subject_context,
    }
    montar = quiz_prompts.prompt_prova_concurso if for_exam else quiz_prompts.prompt_questionario_estudo
    return montar(ctx, avoid_block)


def generate_subjective_quiz_prompt(subject_context, prior_questions=None, question_count=25):
    avoid_block = _avoid_repetition_block(prior_questions)
    n = question_count or 25
    prompt = f"""
Você é um especialista em criar avaliações de nível universitário para a área da saúde.
Analise o texto médico a seguir e elabore um questionário com exatamente {n} questões discursivas (subjetivas).

**INSTRUÇÕES CRÍTICAS DE FORMATAÇÃO E CONTEÚDO:**
- **Formato de Saída:** Sua resposta DEVE ser um único bloco de código JSON sintaticamente perfeito.
- **Estrutura Principal:** O JSON deve conter uma chave principal `"questionario_subjetivo"` cujo valor é um **ARRAY** de objetos.
- **Estrutura da Questão:** Cada objeto deve conter: `pergunta`, `resposta_esperada`, `dificuldade` e `topico`.
- **Qualidade da Resposta Esperada:** O campo `resposta_esperada` deve ser um gabarito completo e detalhado.

**CAMPOS OBRIGATÓRIOS:**
- `pergunta`: O enunciado da questão discursiva
- `resposta_esperada`: Gabarito completo com os pontos-chave esperados na resposta
- `dificuldade`: Nível de dificuldade - usar APENAS: "facil", "medio" ou "dificil"
- `topico`: Categoria temática — escolha a que MELHOR reflete o assunto da questão. **Na dúvida, prefira uma categoria CLÍNICA.** As de saúde coletiva/gestão são SÓ para questões cujo assunto É, de fato, administração de serviços/sistemas, política, financiamento ou saúde coletiva. ⚠️ "gestao_em_saude" significa ADMINISTRAÇÃO de serviços/sistemas de saúde — NÃO conduta do paciente: protocolo/conduta/manejo clínico (ex.: tempo de jejum pré-operatório, ajuste de dose, preparo para exame, profilaxia, indicação cirúrgica) é CLÍNICO → use "tratamento"/"prevencao"/"diagnostico"/etc., NUNCA "gestao_em_saude". Não force categoria clínica em conteúdo de gestão/política — nem o contrário. Clínicas: "fisiopatologia", "quadro_clinico", "diagnostico", "tratamento", "farmacologia", "epidemiologia", "prevencao", "anatomia", "fisiologia", "semiologia". Saúde coletiva/gestão: "saude_coletiva", "politicas_de_saude", "atencao_primaria", "gestao_em_saude", "financiamento_saude", "vigilancia_em_saude", "promocao_da_saude", "etica_e_legislacao"

**REGRAS DE CONTEÚDO:**
- As perguntas devem ser desafiadoras e clinicamente relevantes.
- Distribua as questões entre os níveis de dificuldade: aproximadamente 30% fáceis, 50% médias, 20% difíceis.
- **Cubra TODA a amplitude do material e VARIE entre provas:** distribua as questões por todos os tópicos/seções do conteúdo — não concentre nos mesmos pontos "clássicos" de maior destaque. Explore também detalhes, exceções, comparações, valores/critérios e cenários menos óbvios; a cada nova prova, priorize ângulos e subtemas ainda pouco explorados.

**EXEMPLO DE ESTRUTURA JSON VÁLIDA:**
```json
{{
  "questionario_subjetivo": [
    {{
      "pergunta": "Explique a fisiopatologia da insuficiência cardíaca congestiva e os mecanismos compensatórios iniciais.",
      "resposta_esperada": "Uma resposta completa deve incluir: 1. Definição de IC como síndrome de baixo débito cardíaco; 2. Mecanismos compensatórios: ativação do sistema nervoso simpático, sistema renina-angiotensina-aldosterona, liberação de peptídeos natriuréticos; 3. Remodelamento cardíaco...",
      "dificuldade": "medio",
      "topico": "fisiopatologia"
    }},
    {{
      "pergunta": "Descreva os critérios diagnósticos para diabetes mellitus tipo 2.",
      "resposta_esperada": "Os critérios diagnósticos incluem: 1. Glicemia de jejum ≥126 mg/dL; 2. Glicemia 2h após TOTG ≥200 mg/dL; 3. HbA1c ≥6,5%; 4. Glicemia aleatória ≥200 mg/dL com sintomas clássicos...",
      "dificuldade": "facil",
      "topico": "diagnostico"
    }}
  ]
}}
```

{avoid_block}
**⚠️ REVISÃO DE CONSISTÊNCIA (faça em CADA questão antes de incluí-la):** o enunciado e a resposta esperada precisam ser COERENTES entre si, sem contradição interna. Escreva em português com TODA a acentuação e os cedilhas ("internação", "regência") — jamais em texto sem acento. O texto é renderizado como Markdown inline: use `**negrito**` para destacar os termos que a questão manda analisar (e só cite "termo destacado" se ele estiver mesmo marcado assim), `______` para lacunas, e nada de títulos, listas, tabelas ou HTML. Não cite grifo, cor, figura, tabela ou "linha 3": esses recursos não existem aqui. Em itens de LÍNGUA PORTUGUESA (ou com lacunas), aplique a regra a CADA lacuna/termo separadamente e confira gênero, número, regência e crase antes de fechar a resposta. Atenção REDOBRADA a pares que se invertem e mudam tudo: obstrução **ALTA = delgado** × **BAIXA = cólon**; proximal × distal; aumenta × diminui; sensibilidade × especificidade; causa × efeito; agonista × antagonista. Não diga, por exemplo, "obstrução baixa" tratando de delgado (delgado é obstrução ALTA). Ao citar sigla/mnemônico, confira que cada letra bate com a palavra que representa (Conselho = "Co", não "Cu") e CORRIJA erro óbvio de grafia da fonte (ex.: transcrição "Fucuplarecoco" → "Fucoplarecoco"). Achou incoerência? Corrija antes de finalizar.

Gere exatamente {n} questões subjetivas seguindo esta estrutura JSON.

TEXTO PARA ANÁLISE:
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt

def generate_summary_prompt(subject_context):
    prompt = f"""Você vai criar um resumo de estudo médico no estilo de um estudante de medicina muito organizado e metódico.

REGRAS CRÍTICAS DE ESTILO:
1. NUNCA comece com frases como "Este resumo...", "Esse é um resumo que...", "O presente resumo...", "Neste material..." ou qualquer frase que descreva o que o documento é. Vá DIRETO ao conteúdo.
2. Comece SEMPRE com o título principal do tema (usando ##) seguido imediatamente do conteúdo relevante.
3. Escreva como se fossem suas anotações pessoais de estudo - direto, objetivo e profissional.
4. O tom deve ser técnico-didático, como um estudante explicando para si mesmo ou para colegas.

ESTRUTURA OBRIGATÓRIA:
- Use ## para o título principal do tema
- Use ### para subtópicos e seções
- Use **negrito** para termos-chave, conceitos importantes e palavras que precisam de destaque
- Use listas com marcadores (-) para enumerar características, sintomas, tratamentos
- Use listas numeradas (1.) para sequências, etapas de diagnóstico ou classificações hierárquicas
- Separe seções logicamente com quebras de linha

ORGANIZAÇÃO DO CONTEÚDO:
1. Comece com a definição/conceito central
2. Siga com epidemiologia ou relevância clínica (quando aplicável)
3. Aborde fisiopatologia ou mecanismos
4. Liste manifestações clínicas ou características principais
5. Apresente diagnóstico e/ou diagnósticos diferenciais
6. Finalize com tratamento ou manejo clínico
7. Adicione "Pontos-Chave" ou "Lembre-se" ao final com os conceitos mais importantes para memorização

EXEMPLO DE INÍCIO CORRETO:
## Insuficiência Cardíaca

A **insuficiência cardíaca (IC)** é uma síndrome clínica caracterizada pela incapacidade do coração de manter débito cardíaco adequado às necessidades metabólicas do organismo.

### Classificação
- **IC com fração de ejeção reduzida (ICFEr)**: FE ≤ 40%
- **IC com fração de ejeção preservada (ICFEp)**: FE ≥ 50%
...

TEXTO BASE PARA ANÁLISE:
\"\"\"
{subject_context}
\"\"\"

Sua resposta DEVE ser um único bloco de código JSON com a chave "summary" contendo o resumo em formato Markdown:

```json
{{
  "summary": "## Título do Tema\\n\\nConteúdo do resumo em markdown..."
}}
```"""
    return prompt

def generate_flashcards_prompt(theme, num_flashcards, subject_context):
    prompt = f"""
Você é um especialista na criação de ferramentas de estudo eficazes para o nível de graduação.
Crie {num_flashcards} flashcards sobre o tema '{theme}'{f", considerando o seguinte contexto adicional: {subject_context}" if subject_context else ""}.

Sua resposta DEVE ser um único bloco de código JSON, sem nenhum texto antes ou depois.

```json
{{
  "flashcards": [
    {{
      "frente": "Pergunta ou termo chave do flashcard 1.",
      "verso": "Resposta concisa e informativa para o flashcard 1."
    }},
    {{
      "frente": "Pergunta ou termo chave do flashcard 2.",
      "verso": "Resposta concisa e informativa para o flashcard 2."
    }}
  ]
}}
```
Gere exatamente {num_flashcards} flashcards. As perguntas devem focar em conceitos cruciais, definições importantes e fatos essenciais do tema.
"""
    return prompt


def generate_enhanced_flashcards_prompt(subject_context):
    prompt = f"""Você é um especialista em criar flashcards de alta qualidade para estudantes de medicina.

REGRAS CRÍTICAS:
1. Crie entre 12 a 18 flashcards baseados no conteúdo fornecido
2. VARIE os tipos de pergunta para estimular diferentes níveis cognitivos
3. Distribua as dificuldades: ~30% fácil, ~50% médio, ~20% difícil
4. Categorize cada card por área temática
5. Adicione dicas quando a pergunta for complexa
6. Inclua mnemônicos APENAS quando forem genuinamente úteis e conhecidos

TIPOS DE PERGUNTA (varie entre eles):
- **definicao**: "O que é X?" / "Defina X"
- **mecanismo**: "Por que X causa Y?" / "Qual o mecanismo de X?"
- **clinica**: "Paciente com X e Y apresenta..." / "Qual a principal manifestação de X?"
- **diagnostico**: "Qual exame padrão-ouro para X?" / "Como confirmar X?"
- **tratamento**: "Qual o tratamento de primeira linha para X?" / "Como manejar X?"
- **diferencial**: "Como diferenciar X de Y?"

NÍVEIS DE DIFICULDADE:
- **facil**: Conceitos básicos, definições diretas, fatos isolados
- **medio**: Relações causa-efeito, aplicação clínica, integração de conceitos
- **dificil**: Raciocínio complexo, diagnósticos diferenciais, casos atípicos

ESTRUTURA DE CADA FLASHCARD:
```json
{{
  "frente": "Pergunta clara e objetiva",
  "verso": "Resposta concisa mas completa (máx 2-3 frases)",
  "dica": "Pista sutil para ajudar a lembrar (ou null se desnecessário)",
  "categoria": "fisiopatologia|quadro_clinico|diagnostico|tratamento|farmacologia|epidemiologia",
  "dificuldade": "facil|medio|dificil",
  "tipo": "definicao|mecanismo|clinica|diagnostico|tratamento|diferencial",
  "mnemonico": "Mnemônico conhecido se existir (ou null)"
}}
```

EXEMPLO DE FLASHCARDS BEM ELABORADOS:
```json
{{
  "flashcards": [
    {{
      "frente": "Quais são os critérios de Framingham para IC?",
      "verso": "Maiores: DPN, estase jugular, estertores, cardiomegalia, EAP, B3. Menores: edema MMII, tosse noturna, dispneia aos esforços, hepatomegalia, derrame pleural.",
      "dica": "Pense no que você VÊ no exame físico (maiores) vs. sintomas relatados (menores)",
      "categoria": "diagnostico",
      "dificuldade": "medio",
      "tipo": "diagnostico",
      "mnemonico": null
    }},
    {{
      "frente": "Por que os IECA causam tosse seca?",
      "verso": "Inibem a degradação de bradicinina, que se acumula nas vias aéreas e estimula fibras C sensoriais, causando tosse.",
      "dica": "Lembre que a ECA degrada mais do que só angiotensina...",
      "categoria": "farmacologia",
      "dificuldade": "medio",
      "tipo": "mecanismo",
      "mnemonico": null
    }},
    {{
      "frente": "Qual a tríade clássica da estenose aórtica?",
      "verso": "Angina, síncope e dispneia (insuficiência cardíaca).",
      "dica": "ASD - os sintomas aparecem nessa ordem de progressão",
      "categoria": "quadro_clinico",
      "dificuldade": "facil",
      "tipo": "clinica",
      "mnemonico": "ASD: Angina → Síncope → Dispneia"
    }}
  ]
}}
```

TEXTO BASE PARA CRIAR OS FLASHCARDS:
\"\"\"
{subject_context}
\"\"\"

Sua resposta DEVE ser APENAS o bloco JSON, sem texto adicional."""
    return prompt


def generate_mind_map_prompt(theme, subject_context):
    prompt = f"""
Você é um especialista em criar mapas mentais para estudo médico, focando em clareza, síntese e organização visual.
Sua tarefa é analisar o texto médico fornecido e criar um mapa mental CONCISO e FOCADO nos conceitos mais importantes.

**REGRA CRÍTICA DE PRIORIZAÇÃO:**
- Independente do tamanho do texto, você DEVE selecionar APENAS os 5-6 conceitos MAIS IMPORTANTES clinicamente
- PRIORIZE: definições essenciais, fisiopatologia central, diagnóstico-chave, tratamento de primeira linha
- IGNORE: detalhes secundários, exemplos extensos, informações repetitivas, referências bibliográficas
- O mapa mental deve ser VISUAL e MEMORÁVEL, não um resumo completo do texto

**Instruções de Estrutura:**

1. **Tema Central:** Identifique o tema principal do texto para o campo `tema_central`. Adicione uma `descricao` breve (1 frase) que contextualize o tema.

2. **Ramos Principais (EXATAMENTE 5-6 nós, não mais):** Cada ramo DEVE ter:
   - `titulo`: Frase curta e impactante (máximo 5 palavras)
   - `categoria`: Uma das seguintes: "definicao", "fisiopatologia", "quadro_clinico", "diagnostico", "tratamento", "prognostico", "prevencao", "epidemiologia"
   - `importancia`: "alta", "media" ou "baixa" (baseado na relevância clínica)
   - `descricao`: Frase explicativa concisa (máximo 20 palavras)
   - `filhos`: Sub-nós (máximo 3-4 por ramo)

3. **Hierarquia:** Profundidade máxima de 2 níveis. Filhos devem ter apenas `titulo` e opcionalmente `descricao` curta.

4. **Foco Visual:** O mapa será convertido em imagem - títulos curtos e impactantes são essenciais para legibilidade.

Sua resposta DEVE ser um único bloco de código JSON válido, sem texto adicional.

**Exemplo de Estrutura:**
```json
{{
  "mapa_mental": {{
    "tema_central": "Insuficiência Cardíaca Congestiva",
    "descricao": "Síndrome clínica complexa resultante de alterações estruturais ou funcionais cardíacas",
    "nos": [
      {{
        "titulo": "Fisiopatologia e Mecanismos",
        "categoria": "fisiopatologia",
        "importancia": "alta",
        "descricao": "Compreender os mecanismos é essencial para entender a progressão da doença e o racional terapêutico",
        "filhos": [
          {{"titulo": "Disfunção Sistólica vs Diastólica", "descricao": "FE reduzida (<40%) ou preservada (>50%)"}},
          {{"titulo": "Remodelamento Ventricular"}},
          {{"titulo": "Ativação Neuro-hormonal", "descricao": "SRAA, sistema simpático, peptídeos natriuréticos"}}
        ]
      }},
      {{
        "titulo": "Manifestações Clínicas",
        "categoria": "quadro_clinico",
        "importancia": "alta",
        "descricao": "Sintomas e sinais que guiam o diagnóstico clínico e classificação funcional",
        "filhos": [
          {{"titulo": "Dispneia e Ortopneia"}},
          {{"titulo": "Edema Periférico"}},
          {{"titulo": "Classificação NYHA", "descricao": "Classes I-IV baseadas na limitação funcional"}}
        ]
      }},
      {{
        "titulo": "Abordagem Diagnóstica",
        "categoria": "diagnostico",
        "importancia": "alta",
        "descricao": "Exames complementares para confirmar diagnóstico e identificar etiologia",
        "filhos": [
          {{"titulo": "BNP/NT-proBNP", "importancia": "alta", "descricao": "Biomarcadores com alto valor preditivo negativo"}},
          {{"titulo": "Ecocardiograma", "importancia": "alta"}},
          {{"titulo": "Critérios de Framingham"}}
        ]
      }},
      {{
        "titulo": "Tratamento Farmacológico",
        "categoria": "tratamento",
        "importancia": "alta",
        "descricao": "Medicações que modificam a história natural da doença e melhoram sobrevida",
        "filhos": [
          {{"titulo": "IECA/BRA", "importancia": "alta"}},
          {{"titulo": "Betabloqueadores", "importancia": "alta"}},
          {{"titulo": "Antagonistas da Aldosterona"}},
          {{"titulo": "Inibidores SGLT2", "descricao": "Nova classe com benefício cardiovascular comprovado"}}
        ]
      }},
      {{
        "titulo": "Prognóstico e Seguimento",
        "categoria": "prognostico",
        "importancia": "media",
        "descricao": "Fatores que influenciam a evolução e estratégias de acompanhamento",
        "filhos": [
          {{"titulo": "Fatores de Mau Prognóstico"}},
          {{"titulo": "Monitorização Ambulatorial"}}
        ]
      }}
    ]
  }}
}}
```

**TEXTO PARA ANÁLISE:**
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt


def generate_mind_map_image(outline_data: dict) -> Optional[str]:
    """
    Generates a mind map image using Gemini 2.5 Flash Image based on the outline data.

    Args:
        outline_data: The structured mind map data (mapa_mental dict)

    Returns:
        Base64 encoded image string, or None if generation fails
    """
    try:
        from google import genai
        from google.genai import types
        import base64

        if not Config.GEMINI_API_KEY:
            logger.error("[MIND_MAP_IMAGE] GEMINI_API_KEY não configurada")
            return None

        # Use Pro model for high-quality image generation, fallback to standard model
        image_model = Config.IMAGE_GEN_MODEL_PRO or Config.IMAGE_GEN_MODEL
        if not image_model:
            logger.error("[MIND_MAP_IMAGE] IMAGE_GEN_MODEL_PRO e IMAGE_GEN_MODEL não configuradas")
            return None

        # Build structured text from outline
        central_theme = outline_data.get('tema_central', 'Mapa Mental')
        description = outline_data.get('descricao', '')
        nodes = outline_data.get('nos', [])

        # Create a structured prompt for the mind map
        branches_text = []
        for i, node in enumerate(nodes[:6], 1):  # Limit to 6 main branches for clarity
            category = node.get('categoria', '')
            title = node.get('titulo', '')
            children = node.get('filhos', [])

            branch_str = f"Ramo {i}: {title}"
            if children:
                children_titles = [c.get('titulo', '') for c in children[:4]]  # Limit children
                branch_str += f" (subtópicos: {', '.join(children_titles)})"
            branches_text.append(branch_str)

        branches_formatted = '\n'.join(branches_text)

        # Craft the image generation prompt
        image_prompt = f"""Create a professional medical mind map diagram with the following structure:

CENTRAL TOPIC: "{central_theme}"

MAIN BRANCHES:
{branches_formatted}

STYLE REQUIREMENTS:
- Clean, professional medical/educational style
- Central topic in a prominent oval or rounded rectangle in the center
- Main branches radiating outward with clear connecting lines
- Use a cohesive color scheme: purple for central, then blue, green, teal, orange for branches
- Each branch should be in a rounded rectangle with readable text
- Include subtle icons or visual cues for medical context
- White or very light background for clarity
- Sans-serif font, bold for titles
- Hierarchical layout with clear visual hierarchy
- All text must be clearly legible
- Modern, minimalist aesthetic suitable for medical education
- Aspect ratio approximately 16:9 (landscape)"""

        logger.info(f"[MIND_MAP_IMAGE] Gerando imagem com {image_model}")
        logger.debug(f"[MIND_MAP_IMAGE] Prompt: {image_prompt[:200]}...")

        # Initialize client
        image_client = genai.Client(api_key=Config.GEMINI_API_KEY)

        # Retry logic
        max_retries = 3
        last_error = None

        for attempt in range(max_retries):
            try:
                response = image_client.models.generate_content(
                    model=image_model,
                    contents=image_prompt,
                    config=types.GenerateContentConfig(
                        response_modalities=["IMAGE", "TEXT"],
                    )
                )

                # Extract image from response
                image_data = None
                for part in response.candidates[0].content.parts:
                    if hasattr(part, 'inline_data') and part.inline_data:
                        image_data = part.inline_data.data
                        break

                if image_data:
                    logger.info(f"[MIND_MAP_IMAGE] Imagem gerada com sucesso na tentativa {attempt + 1}")
                    # Return base64 encoded image
                    if isinstance(image_data, bytes):
                        return base64.b64encode(image_data).decode('utf-8')
                    return image_data
                else:
                    # Check for text response
                    text_parts = [p.text for p in response.candidates[0].content.parts if hasattr(p, 'text') and p.text]
                    if text_parts:
                        logger.warning(f"[MIND_MAP_IMAGE] API retornou texto ao invés de imagem: {text_parts[0][:100]}")
                    last_error = "Modelo não gerou imagem"

            except Exception as e:
                logger.warning(f"[MIND_MAP_IMAGE] Tentativa {attempt + 1} falhou: {str(e)}")
                last_error = str(e)

            # Wait before retry
            if attempt < max_retries - 1:
                time.sleep(1 * (attempt + 1))

        logger.error(f"[MIND_MAP_IMAGE] Falha após {max_retries} tentativas: {last_error}")
        return None

    except Exception as e:
        logger.error(f"[MIND_MAP_IMAGE] Erro ao gerar imagem: {str(e)}", exc_info=True)
        return None


def generate_comparative_table_prompt(subject_context):
    prompt = f"""
Você é um especialista em criar materiais de estudo para a área da saúde.
Analise o conteúdo fornecido e crie uma TABELA COMPARATIVA estruturada e enriquecida.

INSTRUÇÕES:
1. Identifique os principais elementos que podem ser comparados (fármacos, doenças, procedimentos, etc.)
2. Crie colunas relevantes para comparação (mecanismo, indicações, contraindicações, efeitos adversos, etc.)
3. Preencha cada célula com informações concisas e clinicamente relevantes
4. Use no máximo 6 colunas e 12 linhas para manter legibilidade
5. Agrupe as linhas por categoria quando apropriado (ex: "1ª linha", "2ª linha", "Reserva")
6. Identifique células com informações críticas (contraindicações absolutas, alertas importantes)
7. Liste todas as abreviações usadas com suas definições

ESTRUTURA JSON OBRIGATÓRIA:

```json
{{
  "table": {{
    "title": "Título descritivo da tabela",
    "headers": ["Coluna1", "Coluna2", "Coluna3", ...],
    "rows": [
      {{
        "cells": ["Valor1", "Valor2", "Valor3", ...],
        "category": "Categoria opcional (ex: 1ª linha, Classe A)",
        "highlights": [2, 4]
      }}
    ],
    "categories": ["1ª linha", "2ª linha"],
    "abbreviations": {{
      "HAS": "Hipertensão Arterial Sistêmica",
      "IC": "Insuficiência Cardíaca"
    }},
    "footnotes": ["Nota de rodapé opcional 1", "Nota 2"]
  }}
}}
```

DETALHES DOS CAMPOS:
- rows: Array de objetos com:
  - cells: Array de strings com os valores de cada coluna
  - category: String opcional para agrupar linhas similares
  - highlights: Array opcional de índices (0-based) das células que contêm informações críticas (contraindicações, alertas)
- categories: Array de strings com as categorias únicas usadas (em ordem de exibição)
- abbreviations: Objeto com siglas como chaves e definições como valores

IMPORTANTE:
- Se não houver categorização natural, omita "category" das rows e "categories" do objeto principal
- Use "highlights" apenas para células com informações críticas de segurança clínica
- Inclua TODAS as abreviações médicas usadas no objeto "abbreviations"

CONTEÚDO PARA ANÁLISE:
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt

def generate_podcast_script_prompt(theme, subject_context):
    prompt = f"""
Você é um roteirista de podcasts médicos educacionais. Sua tarefa é criar um roteiro de diálogo entre dois apresentadores sobre o tema '{theme}', baseado no contexto fornecido.
PERSONAGENS:
Dr. Qython: Um médico especialista, experiente e com uma didática clara. Ele explica os conceitos complexos.
Dra. Epione: A co-apresentadora, que representa o estudante ou profissional em formação. Ela faz as perguntas, pede esclarecimentos e resume os pontos-chave.
ESTRUTURA DO ROTEIRO:
Duração Alvo: O roteiro deve ser dimensionado para um podcast de aproximadamente 5 a 7 minutos (cerca de 800-1000 palavras).
Formato do Diálogo: Cada fala DEVE começar com o nome do personagem, seguido por dois pontos. Ex: Dr. Qython: ou Dr. Epione:.
NÃO inclua sua resposta em um bloco de código JSON. Apenas o texto do roteiro.
EXEMPLO DE SAÍDA ESPERADA:
Dr. Qython: Olá, Dra. Epione! Hoje vamos falar sobre a avaliação de risco cardiovascular em cirurgias não cardíacas. É um tema crucial.
Dra. Epione: Exato, Dr. Qython. Por onde começamos? Quais são os primeiros passos nessa avaliação?
Dr. Qython: O primeiro passo é sempre uma boa anamnese e exame físico, focando em identificar o que chamamos de 'condições cardíacas ativas'...
O diálogo deve ser natural, informativo e seguir uma progressão lógica.
CONTEXTO PARA O ROTEIRO:
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt

def generate_slideshow_only_prompt(subject_context):
    """
    Generates the prompt for creating a rich-content slideshow with explicit
    tags for the context to be analyzed.
    """
    prompt = f"""
Você é um especialista em design instrucional para a área da saúde. Sua tarefa é analisar o "TEXTO PARA ANÁLISE" e criar uma apresentação de slides didática, completa e visualmente rica.
INSTRUÇÕES CRÍTICAS DE FORMATAÇÃO E CONTEÚDO:
SAÍDA OBRIGATÓRIA: Sua resposta DEVE ser um único bloco de código JSON, sintaticamente perfeito. Não inclua nenhum texto, comentário ou json antes ou depois do objeto JSON.
ESTRUTURA PRINCIPAL: O JSON deve ter as chaves title (string), theme (string), e slides (lista de objetos de slide).
ESTRUTURA DO SLIDE: Cada objeto na lista slides deve ter uma chave title (string) e uma chave content (lista de "blocos de conteúdo").
BLOCOS DE CONTEÚDO: A chave content de cada slide deve ser uma lista de objetos. Cada objeto DEVE ter uma chave type que define seu formato. Os tipos permitidos são: text, table, image_suggestion, key_takeaway, e clinical_vignette.
CLAREZA DO SLIDE: Para manter a legibilidade, um slide não deve conter mais do que 2 ou 3 blocos de conteúdo. Se um tópico for muito denso, divida-o em múltiplos slides, adicionando (Parte 1), (Parte 2), etc., aos títulos.
SLIDE DE REFERÊNCIAS OBRIGATÓRIO: O último slide da apresentação DEVE ser dedicado exclusivamente às referências bibliográficas ou fontes de informação utilizadas. O título deste slide DEVE ser "Referências Bibliográficas".
DEFINIÇÃO DOS TIPOS DE BLOCOS DE CONTEÚDO:
"type": "text": Para texto e listas.
"points": Uma lista de strings.
Para criar um subtópico com recuo, inicie a string com dois espaços. Ex: " Este é um subtópico."
"type": "table": Para dados comparativos, classificações ou informações estruturadas.
"title": Um título descritivo para a tabela.
"columns": Uma lista de strings para os cabeçalhos das colunas.
"rows": Uma lista de listas, onde cada lista interna representa uma linha da tabela.
"type": "image_suggestion"**: Para ilustrar um conceito, anatomia ou procedimento.
"description": Uma descrição clara e concisa para uma busca em uma API de imagens. Seja específico. NÃO FORNEÇA URLs. Ex: "Eletrocardiograma mostrando supradesnivelamento do segmento ST em derivações inferiores (DII, DIII, aVF)."
"type": "key_takeaway"**: Para um slide de resumo no final de uma seção importante.
"points": Uma lista de 2 a 4 pontos cruciais e de alto rendimento.
"type": "clinical_vignette"**: Para aplicar o conhecimento a um caso prático.
"title": Título do caso. Ex: "Caso Clínico: Dor Torácica Súbita".
"scenario": A descrição do caso clínico.
"question": Uma pergunta sobre o caso.
"answer": A resposta ou conduta esperada.
EXEMPLO DE ESTRUTURA JSON COMPLETA E VÁLIDA (SEU GUIA):

{{
  "title": "Infarto Agudo do Miocárdio com Supradesnivelamento do Segmento ST (IAMCSST)",
  "theme": "Diagnóstico e Manejo de Emergência",
  "slides": [
    {{
      "title": "Definição e Fisiopatologia",
      "content": [
        {{
          "type": "text",
          "points": [
            "Necrose miocárdica resultante de oclusão coronariana aguda e completa.",
            "Principal causa: Trombose sobre uma placa aterosclerótica rota."
          ]
        }},
        {{
          "type": "image_suggestion",
          "description": "Diagrama de uma artéria coronária com uma placa aterosclerótica rota e formação de trombo oclusivo."
        }}
      ]
    }}
  ]
}}
Agora, analise o texto fornecido dentro das tags <TEXTO_PARA_ANALISE> e gere uma apresentação completa seguindo exatamente a estrutura e as regras descritas.
<TEXTO_PARA_ANALISE>
{subject_context}
</TEXTO_PARA_ANALISE>
"""
    return prompt

def generate_detailed_text_prompt(subject_context):
    prompt = f"""Você é um professor de medicina experiente, reconhecido pela didática excepcional e capacidade de transformar conteúdos complexos em aulas envolventes. Sua tarefa é criar uma AULA COMPLETA baseada no texto fornecido.

REGRAS CRÍTICAS DE ESTILO:
1. NUNCA comece com frases como "Esta aula aborda...", "Neste documento...", "Vamos estudar...", "O objetivo desta aula..." ou qualquer meta-descrição. Vá DIRETO ao conteúdo.
2. Comece SEMPRE com o título principal (usando #) seguido de uma introdução que contextualiza clinicamente o tema.
3. Escreva como um professor experiente explicando para residentes e internos - tom didático, mas profundo.
4. NÃO USE JSON. Sua resposta deve ser APENAS o conteúdo da aula em Markdown puro.

ESTRUTURA OBRIGATÓRIA DA AULA:
Use # para o título principal, ## para seções, ### para subseções.

1. **# Título do Tema**
2. **## Introdução e Relevância Clínica** - Por que este tema importa? Contexto epidemiológico.
3. **## Fisiopatologia** - Mecanismos explicados de forma didática, conectando causa e efeito.
4. **## Quadro Clínico** - Manifestações organizadas (sinais, sintomas, formas de apresentação).
5. **## Diagnóstico** - Anamnese, exame físico, exames complementares, critérios diagnósticos.
6. **## Diagnósticos Diferenciais** - O que pode mimetizar esta condição?
7. **## Tratamento** - Abordagem terapêutica completa (não-farmacológico, farmacológico, cirúrgico se aplicável).
8. **## Prognóstico e Complicações** - Evolução esperada e o que vigiar.
9. **## Pontos para Revisão** - Lista com os 5-8 conceitos mais importantes para memorizar.

ELEMENTOS DIDÁTICOS OBRIGATÓRIOS:
- Use **negrito** para termos-chave e conceitos fundamentais
- Use listas (-) para enumerar características, sintomas, medicamentos
- Use listas numeradas (1.) para sequências e algoritmos
- Inclua pelo menos 2 blocos de destaque usando > (blockquote) para:
  - **> 💡 Pérola Clínica:** dicas práticas que diferenciam um bom médico
  - **> ⚠️ Atenção:** armadilhas diagnósticas ou erros comuns
- Quando relevante, inclua exemplos de casos clínicos curtos para ilustrar conceitos
- Conecte sempre a teoria à prática: "Na prática, isso significa que..."

EXEMPLO DE INÍCIO CORRETO:
# Insuficiência Cardíaca

## Introdução e Relevância Clínica

A **insuficiência cardíaca (IC)** representa uma das principais causas de internação hospitalar em pacientes acima de 65 anos, com prevalência crescente devido ao envelhecimento populacional. Compreender seus mecanismos é fundamental para o manejo adequado destes pacientes.

> 💡 **Pérola Clínica:** Todo paciente com IC descompensada deve ter sua fração de ejeção avaliada - isso muda completamente a estratégia terapêutica.

## Fisiopatologia

O coração insuficiente não consegue manter o **débito cardíaco** adequado às demandas metabólicas do organismo. Isso desencadeia uma cascata de mecanismos compensatórios...

---

TEXTO BASE PARA TRANSFORMAR EM AULA:
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt


def generate_critical_appraisal_prompt(subject_context):
    """Leitura crítica / appraisal de artigo científico (MBE) — estritamente grounded no conteúdo."""
    prompt = f"""
Você é um especialista em Medicina Baseada em Evidências (MBE) e leitura crítica de artigos
científicos, no nível de prova de residência médica.

Faça a LEITURA CRÍTICA do conteúdo científico fornecido. Seja rigoroso, estruturado e didático.

REGRAS CRÍTICAS (não negociáveis):
- Baseie-se ESTRITAMENTE no conteúdo fornecido. NÃO invente dados, números, intervalos de
  confiança, p-valores, referências ou desfechos que não estejam no material.
- Se algo não estiver no material, escreva "não informado no material" — nunca chute.
- Não fabrique citações nem resultados.

INSTRUÇÕES:
1. Identifique o tipo de estudo (ECR, coorte, caso-controle, transversal, revisão sistemática,
   meta-análise, relato/série de casos, etc.).
2. Extraia o PICO (População, Intervenção, Comparação, Desfecho).
3. Avalie o risco de viés por domínio (randomização, alocação/sigilo, cegamento, perdas de
   seguimento, relato seletivo) com julgamento (baixo/incerto/alto) e justificativa.
4. Classifique o nível de evidência (Oxford CEBM) e a força/certeza (GRADE).
5. Resuma os resultados-chave com a medida de efeito quando houver (RR, OR, HR, NNT, IC95%, p).
6. Liste forças e limitações.
7. Avalie a aplicabilidade à prática (inclusive realidade brasileira/SUS quando pertinente).
8. Diga COMO esse tema/desenho costuma ser cobrado em prova de residência.
9. Dê o "bottom line" (conclusão prática em 1-2 frases).

Sua resposta DEVE ser um único bloco JSON, sem texto fora dele:

```json
{{
  "appraisal": {{
    "title": "Leitura Crítica: [tema/título do estudo]",
    "citation": "referência se identificável no material, senão 'não informado no material'",
    "study_type": "tipo de estudo",
    "objective": "objetivo do estudo em 1 frase",
    "pico": {{
      "population": "...",
      "intervention": "...",
      "comparison": "...",
      "outcome": "..."
    }},
    "evidence": {{
      "oxford_level": "ex.: 1b (ou 'não informado no material')",
      "grade": "Alta | Moderada | Baixa | Muito baixa",
      "rationale": "por que esse nível/certeza"
    }},
    "risk_of_bias": [
      {{ "domain": "Randomização", "judgment": "baixo | incerto | alto", "rationale": "..." }}
    ],
    "key_results": [
      {{ "outcome": "...", "effect": "ex.: RR 0,75 (IC95% 0,60-0,94); p=0,01 — ou 'não informado no material'", "interpretation": "significado clínico" }}
    ],
    "strengths": ["..."],
    "limitations": ["..."],
    "applicability": "aplicabilidade à prática (contexto BR/SUS quando couber)",
    "exam_relevance": "como costuma cair em prova de residência",
    "bottom_line": "conclusão prática em 1-2 frases"
  }}
}}
```

CONTEÚDO PARA ANÁLISE:
{subject_context}
"""
    return prompt


def generate_clinical_case_prompt(subject_context):
    """Generate an interactive clinical case with branching decisions."""
    prompt = f"""
Você é um especialista em educação médica baseada em casos clínicos.
Crie um CASO CLÍNICO INTERATIVO com decisões ramificadas baseado no conteúdo fornecido.

INSTRUÇÕES:
1. Crie um cenário clínico realístico e educativo
2. Inclua 4-6 pontos de decisão
3. Cada decisão deve ter 2-3 opções com feedback
4. Uma opção deve ser ótima, outras subótimas (todas educativas)
5. Decisões levam a diferentes desfechos
6. Inclua dados clínicos realistas (sinais vitais, exames)

Sua resposta DEVE ser um único bloco JSON:

```json
{{
  "clinical_case": {{
    "title": "Título do Caso Clínico",
    "difficulty": "intermediate",
    "speciality": "Especialidade médica",
    "patient": {{
      "age": 55,
      "gender": "M",
      "complaint": "Queixa principal do paciente"
    }},
    "blocks": [
      {{
        "id": "start",
        "content": "Descrição inicial do caso clínico. Paciente chega ao PS com...",
        "vitals": {{"PA": "150/95", "FC": "92", "FR": "18", "SpO2": "96%", "Temp": "36.8°C"}},
        "decision": {{
          "question": "Qual sua primeira conduta?",
          "options": [
            {{
              "id": "option_a",
              "text": "Opção de conduta A",
              "next": "block_a",
              "is_best": true,
              "feedback": "Explicação detalhada de por que esta é a melhor opção.",
              "points": 10
            }},
            {{
              "id": "option_b",
              "text": "Opção de conduta B",
              "next": "block_b",
              "is_best": false,
              "feedback": "Explicação de por que esta opção é subótima, mas ainda educativa.",
              "points": 5
            }}
          ]
        }}
      }},
      {{
        "id": "block_a",
        "content": "Evolução do caso após a decisão A...",
        "exam_results": {{"ECG": "Descrição do resultado", "Labs": "Valores laboratoriais"}},
        "decision": {{
          "question": "Próxima pergunta decisória",
          "options": [...]
        }}
      }},
      {{
        "id": "conclusion",
        "type": "end",
        "content": "Resumo do desfecho do caso.",
        "summary": "O paciente evoluiu com...",
        "learning_points": [
          "Ponto de aprendizado 1",
          "Ponto de aprendizado 2",
          "Ponto de aprendizado 3"
        ]
      }}
    ]
  }}
}}
```

CONTEÚDO BASE PARA O CASO:
\"\"\"
{subject_context}
\"\"\"
"""
    return prompt


# --- AUSSIE ENGINEER'S FIX ---
# New prompt and function specifically for generating a video lesson narration script.
def generate_video_narration_prompt(slideshow_json_str: str):
    """
    Generates a prompt to create a narration script based on the structured
    content of a slideshow.
    """
    prompt = f"""
Você é um roteirista de videoaulas médicas. Sua tarefa é criar um roteiro de narração em formato de diálogo entre dois apresentadores, baseado no conteúdo de uma apresentação de slides fornecida em formato JSON.

**PERSONAGENS:**
- **Dr. Qython:** Um médico especialista, experiente e com uma didática clara. Ele explica os conceitos complexos de cada slide.
- **Dra. Epione:** A co-apresentadora, que guia a aula, introduz os slides, faz perguntas pertinentes e resume os pontos-chave para o espectador.

**ESTRUTURA DO ROTEIRO:**
- **Formato do Diálogo:** Cada fala DEVE começar com o nome do personagem, seguido por dois pontos. Ex: `Dr. Qython:` ou `Dra. Epione:`.
- **Fluxo da Aula:** O roteiro DEVE seguir a ordem dos slides da apresentação. A Dra. Epione geralmente introduz o tópico de um novo slide, e o Dr. Qython aprofunda o conteúdo.
- **Sincronia com o Visual:** A narração deve corresponder diretamente ao conteúdo do slide que está sendo apresentado. Faça referências claras ao que está na tela (ex: "Como podemos ver nesta tabela...", "Este diagrama ilustra bem...").
- **Tom:** Profissional, didático e conversacional, como dois professores apresentando uma aula juntos.
- **SAÍDA:** A resposta deve ser APENAS o texto do roteiro, sem JSON, comentários ou qualquer outro texto.

**EXEMPLO DE SAÍDA ESPERADA:**
Dra. Epione: Olá a todos e bem-vindos. No primeiro slide, vamos começar com a definição e a fisiopatologia do Infarto Agudo do Miocárdio. Dr. Qython, pode nos explicar os pontos principais?
Dr. Qython: Com certeza, Dra. Epione. Como vemos nos pontos, o infarto é essencialmente a necrose do músculo cardíaco causada por uma oclusão completa de uma artéria coronária. A causa mais comum é a trombose que se forma sobre uma placa de gordura que se rompeu.
Dra. Epione: E o diagrama na tela ilustra exatamente esse processo, correto?
Dr. Qython: Exato. Vemos a artéria, a placa rompida e o trombo bloqueando o fluxo sanguíneo.

**APRESENTAÇÃO JSON PARA CRIAR O ROTEIRO:**
```json
{slideshow_json_str}
ROTEIRO DA NARRAÇÃO:
"""
    return prompt

def generate_narration_script_from_slideshow(slideshow_json: dict) -> str:
    """
    Uses an LLM to generate a narration script directly from a slideshow's JSON data.
    Includes retry logic with fallback model support.
    """
    logger.info("Generating video lesson narration script from slideshow JSON...")

    # Convert the dictionary to a nicely formatted JSON string for the prompt
    slideshow_str = json.dumps(slideshow_json, indent=2, ensure_ascii=False)
    prompt_text = generate_video_narration_prompt(slideshow_str)

    # Build list of models to try (primary + fallback)
    models_to_try = [PRIMARY_LLM_MODEL]
    if FALLBACK_LLM_MODEL and FALLBACK_LLM_MODEL != PRIMARY_LLM_MODEL:
        models_to_try.append(FALLBACK_LLM_MODEL)
        logger.info(f"[VIDEO_NARRATION] Fallback configurado: primário={PRIMARY_LLM_MODEL}, fallback={FALLBACK_LLM_MODEL}")

    last_error = None
    script = None

    for model_idx, model_to_use in enumerate(models_to_try):
        is_fallback = model_idx > 0

        for attempt in range(MAX_RETRIES):
            try:
                logger.info(f"[VIDEO_NARRATION] Tentativa {attempt + 1}/{MAX_RETRIES} com modelo: {model_to_use}{' (fallback)' if is_fallback else ''}")

                response = client.models.generate_content(
                    model=f'models/{model_to_use}',
                    contents=prompt_text,
                    config=types.GenerateContentConfig(temperature=0.5, max_output_tokens=16384)
                )

                script = response.text.strip() if response.text else None

                if not script:
                    logger.warning(f"[VIDEO_NARRATION] LLM retornou resposta vazia (tentativa {attempt + 1}, modelo: {model_to_use})")
                    last_error = "LLM returned empty response"
                    if attempt < MAX_RETRIES - 1:
                        backoff = min(INITIAL_BACKOFF * (2 ** attempt), MAX_BACKOFF)
                        logger.info(f"[VIDEO_NARRATION] Aguardando {backoff}s antes de retry...")
                        time.sleep(backoff)
                    continue

                # Success
                if is_fallback:
                    logger.warning(f"[VIDEO_NARRATION] Sucesso com modelo de fallback: {model_to_use}")
                logger.info(f"Video narration script generated successfully. Length: {len(script)} chars.")
                return script

            except Exception as e:
                error_msg = str(e)
                last_error = error_msg
                logger.warning(f"[VIDEO_NARRATION] Erro na tentativa {attempt + 1}/{MAX_RETRIES} com modelo {model_to_use}: {error_msg}")

                # Check for 503 (model overloaded) - worth retrying
                if "503" in error_msg or "UNAVAILABLE" in error_msg or "overloaded" in error_msg.lower():
                    if attempt < MAX_RETRIES - 1:
                        backoff = min(INITIAL_BACKOFF * (2 ** attempt), MAX_BACKOFF)
                        logger.info(f"[VIDEO_NARRATION] Erro 503/UNAVAILABLE. Aguardando {backoff}s antes de retry...")
                        time.sleep(backoff)
                    continue
                else:
                    # Other errors - don't retry, try next model
                    break

        # If we got a successful response, we already returned above
        # If not, continue to next model

    # If all models and retries failed, raise the last error
    logger.error(f"[VIDEO_NARRATION] Todos os modelos e tentativas falharam. Último erro: {last_error}")
    raise RuntimeError(f"Failed to generate narration script after all retries: {last_error}")
# --- MAIN SERVICE FUNCTION ---
def generate_study_material(document_content: str, material_type: str, user_id: int, prior_questions=None, question_count=None, banca_profile=None, num_alternatives=5, model_override=None, fallback_override=None, thinking_override=None, for_exam: bool = False) -> Tuple[dict, Any, str]:
    logger.info(f"Generating '{material_type}' material for user_id {user_id}.")
    subject_context = document_content
    prior_questions = prior_questions or []
    is_json_output = True
    prompt_text = ""

    # Step 1: Select the correct prompt based on the material type.
    if material_type == 'slideshow_only':
        prompt_text = generate_slideshow_only_prompt(subject_context)
    elif material_type == 'questionnaire_objective':
        prompt_text = generate_objective_quiz_prompt(subject_context, prior_questions, question_count, num_alternatives, for_exam=for_exam)
        # Now uses JSON output format
    elif material_type == 'podcast':
        prompt_text = generate_podcast_script_prompt("podcast", subject_context)
        is_json_output = False
    elif material_type == 'detailed_text':
        prompt_text = generate_detailed_text_prompt(subject_context)
        is_json_output = False
    elif material_type == 'summary':
        prompt_text = generate_summary_prompt(subject_context)
    elif material_type == 'questionnaire_subjective':
        prompt_text = generate_subjective_quiz_prompt(subject_context, prior_questions, question_count)
    elif material_type == 'flashcards':
        prompt_text = generate_enhanced_flashcards_prompt(subject_context)
    elif material_type == 'mind_map':
        prompt_text = generate_mind_map_prompt("mind_map", subject_context)
    elif material_type == 'comparative_table':
        prompt_text = generate_comparative_table_prompt(subject_context)
    elif material_type == 'clinical_case':
        prompt_text = generate_clinical_case_prompt(subject_context)
    elif material_type == 'critical_appraisal':
        prompt_text = generate_critical_appraisal_prompt(subject_context)
    else:
        logger.error(f"No specific prompt defined for material type '{material_type}'.")
        raise ValueError(f"Unsupported material type: {material_type}")

    # Perfil da banca (dossiê confirmado do card "Meus Concursos"): injeta no topo do prompt,
    # só p/ questionários, para o modelo imitar formato/estilo/temas da prova-alvo.
    if banca_profile and material_type in ('questionnaire_objective', 'questionnaire_subjective'):
        prompt_text = _banca_profile_block(banca_profile) + "\n" + prompt_text

    try:
        # Step 2: Call the LLM to get the raw response with retry and fallback logic.
        # Questionários usam temperatura mais alta p/ variar as questões entre gerações
        # (a fidelidade vem do texto-fonte + da avoid-list de questões já criadas); os
        # demais materiais seguem conservadores.
        quiz_temperature = 0.85 if material_type in ('questionnaire_objective', 'questionnaire_subjective') else 0.4
        genai_config = {'temperature': quiz_temperature, 'max_output_tokens': 32768}
        if is_json_output:
            genai_config['response_mime_type'] = 'application/json'
        # Structured output schema por tipo: constrained decoding garante JSON válido
        # ao nível do modelo, não só na recuperação pós-fato (json-repair é backstop).
        if material_type == 'questionnaire_objective':
            genai_config['response_schema'] = _schema_questionario_objetivo(for_exam)

        # Build list of models to try (primary + fallback)
        # Material usa um modelo DEDICADO (Config.MATERIAL_LLM_MODEL, default lite),
        # isolado do PRIMARY (consultas/resumos/relatórios): material a 3.5-flash dava
        # custo ~6× a receita. Fallback p/ outro modelo barato. Tunável por env sem deploy.
        # model_override: provas de concurso (Meus Concursos) rodam em modelo FORTE, com
        # preço próprio (100 dracmas) — ver Config.EXAM_LLM_MODEL.
        primary_material_model = model_override or Config.MATERIAL_LLM_MODEL or PRIMARY_LLM_MODEL
        fallback_material_model = (
            fallback_override if model_override else (Config.MATERIAL_FALLBACK_LLM_MODEL or FALLBACK_LLM_MODEL)
        )
        models_to_try = [primary_material_model]
        if fallback_material_model and fallback_material_model != primary_material_model:
            models_to_try.append(fallback_material_model)
            logger.info(f"[MATERIAL_GEN] Fallback configurado: primário={primary_material_model}, fallback={fallback_material_model}")

        last_error = None
        raw_output = None
        usage = None
        model_name = None

        for model_idx, model_to_use in enumerate(models_to_try):
            is_fallback = model_idx > 0

            for attempt in range(MAX_RETRIES):
                try:
                    logger.info(f"[MATERIAL_GEN] Tentativa {attempt + 1}/{MAX_RETRIES} com modelo: {model_to_use}{' (fallback)' if is_fallback else ''}")

                    # Thinking POR MODELO: o helper resolve 3.x=thinking_level vs 2.5=thinking_budget
                    # (assim o fallback 2.5-flash-lite não quebra com 400). No lite o thinking é
                    # baratíssimo, então ALTO recupera a profundidade das justificativas.
                    model_config = dict(genai_config)
                    model_config['thinking_config'] = _get_thinking_config_for_model(
                        model_to_use, thinking_override or Config.MATERIAL_THINKING_LEVEL
                    )
                    response = client.models.generate_content(
                        model=f'models/{model_to_use}',
                        contents=prompt_text,
                        config=types.GenerateContentConfig(**model_config)
                    )
                    raw_output = response.text
                    usage = response.usage_metadata
                    model_name = model_to_use

                    if not raw_output:
                        logger.warning(f"[MATERIAL_GEN] LLM retornou resposta vazia (tentativa {attempt + 1}, modelo: {model_to_use})")
                        last_error = "LLM returned empty response"
                        if attempt < MAX_RETRIES - 1:
                            backoff = min(INITIAL_BACKOFF * (2 ** attempt), MAX_BACKOFF)
                            logger.info(f"[MATERIAL_GEN] Aguardando {backoff}s antes de retry...")
                            time.sleep(backoff)
                        continue

                    # Success
                    if is_fallback:
                        logger.warning(f"[MATERIAL_GEN] Sucesso com modelo de fallback: {model_to_use}")
                    break  # Exit retry loop on success

                except Exception as e:
                    error_msg = str(e)
                    last_error = error_msg
                    logger.warning(f"[MATERIAL_GEN] Erro na tentativa {attempt + 1}/{MAX_RETRIES} com modelo {model_to_use}: {error_msg}")

                    # Check for 503 (model overloaded) - worth retrying
                    if "503" in error_msg or "UNAVAILABLE" in error_msg or "overloaded" in error_msg.lower():
                        if attempt < MAX_RETRIES - 1:
                            backoff = min(INITIAL_BACKOFF * (2 ** attempt), MAX_BACKOFF)
                            logger.info(f"[MATERIAL_GEN] Erro 503/UNAVAILABLE. Aguardando {backoff}s antes de retry...")
                            time.sleep(backoff)
                        continue
                    else:
                        # Other errors - don't retry, try next model
                        break

            # If we got a successful response, break out of the model loop
            if raw_output:
                break

        # If all models and retries failed, raise the last error
        if not raw_output:
            raise ValueError(f"LLM returned an empty response for '{material_type}' after all retries. Last error: {last_error}")

        # Instrumentação de custo real (mesma do chat/RAG): material roda no modelo
        # caro (PRIMARY_LLM_MODEL), então logamos in/out/think → $ por material p/
        # medir a margem vs. o preço em dracmas. Defensivo: nunca quebra a geração.
        try:
            from ..llm_services import _log_call_cost
            _log_call_cost(f'material:{material_type}', model_name, usage)
        except Exception:
            pass

        # Step 2b: Para tipos JSON, pré-valida o parse ANTES do Step 3. Se o JSON vier
        # malformado internamente (não só lixo no final, que o raw_decode/json-repair já
        # cobre), tenta o modelo fallback enquanto prompt_text ainda está em escopo.
        # Garante que falhas de parse nunca cheguem ao usuário sem ao menos uma tentativa
        # de recuperação com outro modelo.
        _pre_parsed_json = None
        if is_json_output:
            try:
                _pre_parsed_json = extract_json_from_response(raw_output)
            except ValueError as _parse_err:
                _fb_model = Config.MATERIAL_FALLBACK_LLM_MODEL or FALLBACK_LLM_MODEL
                if _fb_model and _fb_model != model_name:
                    logger.warning(
                        f"[MATERIAL_GEN] JSON inválido após json-repair (model={model_name}) "
                        f"— retentando com {_fb_model}: {_parse_err}"
                    )
                    _fb_cfg = dict(genai_config)
                    _fb_cfg['thinking_config'] = _get_thinking_config_for_model(
                        _fb_model, Config.MATERIAL_THINKING_LEVEL
                    )
                    _fb_resp = client.models.generate_content(
                        model=f'models/{_fb_model}',
                        contents=prompt_text,
                        config=types.GenerateContentConfig(**_fb_cfg)
                    )
                    if _fb_resp and _fb_resp.text:
                        _pre_parsed_json = extract_json_from_response(_fb_resp.text)
                        raw_output = _fb_resp.text
                        model_name = _fb_model
                        usage = _fb_resp.usage_metadata
                        logger.info(f"[MATERIAL_GEN] Fallback {_fb_model} produziu JSON válido")
                    else:
                        raise _parse_err
                else:
                    raise _parse_err

        # Step 3: Process the raw response based on the material type.
        part_dict = {}
        if material_type == 'slideshow_only':
            json_content = _pre_parsed_json
            
            # This list will hold paths for images that need to be embedded in the PPTX
            # and then cleaned up.
            temp_files_for_pptx_cleanup = []
            
            # This is the public-facing directory for previews. We won't clean this up immediately.
            public_preview_dir = os.path.join(Config.PERMANENT_UPLOAD_FOLDER, 'temp_preview_images')
            os.makedirs(public_preview_dir, exist_ok=True)

            for slide in json_content.get("slides", []):
                for content_block in slide.get("content", []):
                    if content_block.get("type") == "image_suggestion":
                        description = content_block.get("description")
                        # Generate the image using our robust service
                        temp_path_from_service = stock_image_service.generate_image_from_prompt(description)
                        
                        if temp_path_from_service:
                            filename = os.path.basename(temp_path_from_service)
                            
                            # 1. Create a permanent copy for the web preview
                            public_path = os.path.join(public_preview_dir, filename)
                            shutil.copy(temp_path_from_service, public_path)
                            
                            # 2. Assign the public URL for the frontend
                            content_block['generated_image_url'] = f"/{Config.STATIC_URL_PATH_PREFIX.strip('/')}/uploads/temp_preview_images/{filename}"
                            
                            # 3. Keep track of the path for PPTX embedding
                            content_block['temp_image_path'] = public_path 
                            temp_files_for_pptx_cleanup.append(public_path)

                            # The original file from the service is no longer needed
                            os.remove(temp_path_from_service)
            
            file_path = create_presentation_from_json(json_content, user_id)
            
            # Now, we clean up the files we copied for the preview/pptx.
            # A better long-term solution is a cron job to clean this folder periodically.
            # For now, we'll leave them so the preview works. The user can close the modal.
            # The key is that the 404 is gone.
            
            # Let's just log that we are NOT cleaning them for now.
            logger.info(f"{len(temp_files_for_pptx_cleanup)} preview images were generated and will remain for the user session.")

            part_dict = {
                "slideshow_file_path": file_path,
                "slideshow_content": json_content
            }

        elif material_type == 'questionnaire_objective':
            part_dict = _pre_parsed_json
            _shuffle_quiz_positions(part_dict)  # fim do viés de posição (correta sempre em B/C/D)
        elif material_type == 'podcast':
            part_dict = {"podcast_script": raw_output}
        elif material_type == 'detailed_text':
            part_dict = {"detailed_text": raw_output}
        elif material_type == 'flashcards':
            part_dict = _pre_parsed_json
        elif material_type == 'mind_map':
            part_dict = _pre_parsed_json

            # Generate the mind map image using the outline data
            if 'mapa_mental' in part_dict:
                logger.info("[MIND_MAP] Gerando imagem do mapa mental...")
                mind_map_image = generate_mind_map_image(part_dict['mapa_mental'])
                if mind_map_image:
                    part_dict['mind_map_image'] = mind_map_image
                    logger.info("[MIND_MAP] Imagem gerada com sucesso")
                else:
                    logger.warning("[MIND_MAP] Falha ao gerar imagem, retornando apenas outline")
                    part_dict['mind_map_image'] = None
            else:
                logger.warning("[MIND_MAP] Estrutura mapa_mental não encontrada no JSON")
                part_dict['mind_map_image'] = None
        else:
            part_dict = _pre_parsed_json if _pre_parsed_json is not None else extract_json_from_response(raw_output)

        return part_dict, usage, model_name

    except Exception as e:
        logger.error(f"Failure in material generation pipeline for '{material_type}': {e}", exc_info=True)
        raise e
